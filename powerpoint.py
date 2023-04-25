from pptx import Presentation
from pptx.util import Inches
import os

def replace_first_image(slide, new_image_path):
    # Find the first picture placeholder in the slide
    for shape in slide.shapes:
        if shape.shape_type == 13:  # ShapeType.PICTURE == 13
            # Get the image properties
            left, top, width, height = shape.left, shape.top, shape.width, shape.height

            # Remove the old image
            slide.shapes._spTree.remove(shape._element)
            
            # Add the new image
            slide.shapes.add_picture(new_image_path, left, top, width, height)
            break

# Load the pre-made PowerPoint file
ppt_file = 'input.pptx'
new_image_path = 'new_image.jpg'
output_file = 'output.pptx'

presentation = Presentation(ppt_file)

# Access the 3rd slide (0-based indexing)
third_slide = presentation.slides[2]

# Replace the first image in the 3rd slide with a new image
replace_first_image(third_slide, new_image_path)

# Save the updated presentation
presentation.save(output_file)
